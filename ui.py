# ui.py

# Standard Python imports
import sys
import os
import time
import re 
import html # For escaping log messages

# PyQt6 imports
from PyQt6.QtWidgets import (QApplication,
                             QMainWindow, QVBoxLayout, QHBoxLayout,
                             QWidget, QLabel, QPushButton, QTextEdit,
                             QLineEdit,
                             QFileDialog, QMessageBox, QSizePolicy,
                             QProgressDialog, QSpacerItem,
                             QFrame,
                             QComboBox,
                             QSpinBox,
                             QStackedWidget)
from PyQt6.QtGui import QFont, QPixmap, QPalette, QColor, QBrush
from PyQt6.QtCore import Qt, QThread, pyqtSignal, QSize, QRect, QTimer 

# Other imports
from transformers import AutoTokenizer
# Assumes worker.py has the updated AIWorker with 3 args in finished signal
from worker import AIWorker 
# Assumes pdf_generator.py and pdf_worker.py have enhanced error reporting
from pdf_generator import generate_pdf 
from pdf_worker import PDFWorker 
from config import * 
from prompts import PREDEFINED_PROMPTS, PROMPT_NAMES

# Import for DOCX handling
try:
    import docx
    from docx.enum.style import WD_STYLE_TYPE 
    from docx.enum.text import WD_UNDERLINE 
except ImportError:
    docx = None
    WD_STYLE_TYPE = None 
    WD_UNDERLINE = None 
    print("Warning: python-docx library not found. DOCX file support will be unavailable.")

# +++ Add python-pptx import +++
try:
    from pptx import Presentation
except ImportError:
    Presentation = None # Assign None if import fails
    print("Warning: python-pptx library not found. PPTX file support will be unavailable.")


# --- Tokenizer Initialization (Lazy Loading) ---
class ModernHackerPDFConverterWindow(QMainWindow):
    def __init__(self):
        super().__init__()
        self.setWindowTitle("FormatAI PDF)")
        self.setGeometry(100, 100, WINDOW_WIDTH, WINDOW_HEIGHT)
        self.setStyleSheet(self._get_stylesheet())
        self._tokenizer_instance = None
        self.token_update_timer = QTimer(self); self.token_update_timer.setSingleShot(True); self.token_update_timer.setInterval(500); self.token_update_timer.timeout.connect(self._perform_token_update)
        central_widget = QWidget(); central_widget.setObjectName("centralWidget"); self.setCentralWidget(central_widget)
        central_layout = QVBoxLayout(central_widget); central_layout.setContentsMargins(LAYOUT_MARGIN, LAYOUT_MARGIN, LAYOUT_MARGIN, LAYOUT_MARGIN); central_layout.setSpacing(0) 
        self.stacked_widget = QStackedWidget(); central_layout.addWidget(self.stacked_widget)
        self._create_status_log_widgets() 
        self.input_page = self._create_input_page()
        self.status_page = self._create_status_page()
        self.stacked_widget.addWidget(self.input_page)
        self.stacked_widget.addWidget(self.status_page)
        self.stacked_widget.setCurrentIndex(0)
        self._apply_background_image()
        self.ai_worker = None; self.pdf_worker = None; self.progress_dialog = None
        self.log_message("Application started.")
        if docx is None: self.log_message("python-docx not found. DOCX loading disabled.", COLOR_WARNING_YELLOW)
        # +++ Added check for Presentation +++
        if Presentation is None: self.log_message("python-pptx not found. PPTX loading disabled.", COLOR_WARNING_YELLOW) 


    def _create_separator(self):
        separator = QFrame(); separator.setFrameShape(QFrame.Shape.HLine); separator.setFrameShadow(QFrame.Shadow.Sunken); return separator
    def _create_status_log_widgets(self):
        self.status_label = QLabel("Status:"); self.status_display = QLineEdit(); self.status_display.setReadOnly(True); self.status_display.setPlaceholderText("Application ready.")
        self.log_label = QLabel("Activity Log:"); self.log_display = QTextEdit(); self.log_display.setReadOnly(True); self.log_display.setMinimumHeight(LOG_AREA_MIN_HEIGHT); self.log_display.setPlaceholderText("Application events..."); self.log_display.setStyleSheet(f"QTextEdit {{ background-color: {COLOR_LOG_BACKGROUND}; color: {COLOR_TEXT_NEON_GREEN}; border: 1px solid {COLOR_HIGHLIGHT_CYAN}; padding: 8px; font-size: 12px; border-radius: 4px; }}")
    def _create_input_page(self):
        page = QWidget(); page_layout = QVBoxLayout(page); page_layout.setContentsMargins(0, 0, 0, 0); page_layout.setSpacing(LAYOUT_SPACING)
        prompt_select_layout = QHBoxLayout(); self.prompt_select_label = QLabel("Choose Prompt Template:"); prompt_select_layout.addWidget(self.prompt_select_label); self.prompt_combo = QComboBox(); self.prompt_combo.addItems(PROMPT_NAMES); self.prompt_combo.currentIndexChanged.connect(self._load_selected_prompt); prompt_select_layout.addWidget(self.prompt_combo); prompt_select_layout.addStretch(1); page_layout.addLayout(prompt_select_layout)
        prompt_layout = QVBoxLayout(); self.prompt_label = QLabel("Editable Prompt Instructions:"); prompt_layout.addWidget(self.prompt_label); self.prompt_input = QTextEdit(); self.prompt_input.setPlaceholderText("Enter instructions..."); self.prompt_input.setMinimumHeight(PROMPT_INPUT_MIN_HEIGHT); prompt_layout.addWidget(self.prompt_input); page_layout.addLayout(prompt_layout)
        page_layout.addWidget(self._create_separator())
        settings_frame = QFrame(); settings_frame.setFrameShape(QFrame.Shape.StyledPanel); settings_frame.setMinimumHeight(SETTINGS_FRAME_HEIGHT); settings_layout = QHBoxLayout(settings_frame); settings_layout.setContentsMargins(10, 5, 10, 5); settings_layout.setSpacing(15); self.page_size_label = QLabel("Page Size:"); settings_layout.addWidget(self.page_size_label); self.page_size_combo = QComboBox(); self.page_size_combo.addItems(PDF_PAGE_SIZE_OPTIONS); self.page_size_combo.setCurrentText(PDF_PAGE_SIZE_DEFAULT); self.page_size_combo.setMinimumWidth(100); settings_layout.addWidget(self.page_size_combo); settings_layout.addSpacerItem(QSpacerItem(20, 20, QSizePolicy.Policy.Fixed, QSizePolicy.Policy.Minimum)); self.font_size_label = QLabel("Font Size:"); settings_layout.addWidget(self.font_size_label); self.font_size_spinbox = QSpinBox();
        if PDF_FONT_SIZE_OPTIONS: self.font_size_spinbox.setRange(min(PDF_FONT_SIZE_OPTIONS), max(PDF_FONT_SIZE_OPTIONS)); self.font_size_spinbox.setValue(PDF_FONT_SIZE_DEFAULT if PDF_FONT_SIZE_DEFAULT in PDF_FONT_SIZE_OPTIONS else min(PDF_FONT_SIZE_OPTIONS))
        else: self.font_size_spinbox.setRange(8, 24); self.font_size_spinbox.setValue(12)
        self.font_size_spinbox.setSingleStep(1); self.font_size_spinbox.setToolTip("Base font size for the PDF."); self.font_size_spinbox.setMinimumWidth(60); settings_layout.addWidget(self.font_size_spinbox); settings_layout.addStretch(1); page_layout.addWidget(settings_frame)
        page_layout.addWidget(self._create_separator())
        original_text_layout = QVBoxLayout(); self.original_text_label = QLabel("Original Text (Paste or Load File):"); original_text_layout.addWidget(self.original_text_label); self.original_text_input = QTextEdit(); self.original_text_input.setPlaceholderText("Paste your original text here..."); self.original_text_input.setMinimumHeight(ORIGINAL_TEXT_INPUT_MIN_HEIGHT); 
        self.original_text_input.textChanged.connect(self._request_token_update) 
        original_text_layout.addWidget(self.original_text_input); token_count_layout = QHBoxLayout(); token_count_layout.addStretch(1); self.token_count_label = QLabel(f"Tokens: 0 / {LLM_CONTEXT_WINDOW}"); self.token_count_label.setStyleSheet(f"color: {COLOR_TOKEN_NORMAL};"); self.token_count_label.setFont(QFont("Consolas", 10)); token_count_layout.addWidget(self.token_count_label); original_text_layout.addLayout(token_count_layout); page_layout.addLayout(original_text_layout) 
        # REMOVED Conversion Mode Section
        page_layout.addSpacerItem(QSpacerItem(20, 10, QSizePolicy.Policy.Minimum, QSizePolicy.Policy.Expanding))
        # --- Bottom Buttons Section --- (Adding Exit Button) ---
        button_container = QWidget()
        button_layout = QHBoxLayout(button_container)
        button_layout.setContentsMargins(0, 0, 0, 0)
        button_layout.setSpacing(10)
        # +++ Add Exit Button +++
        self.exit_button = QPushButton("Exit")
        self.exit_button.setMinimumHeight(BUTTON_MIN_HEIGHT)
        self.exit_button.clicked.connect(self.close) 
        button_layout.addWidget(self.exit_button) 
        button_layout.addStretch(1) # Push other buttons right
        # Load File Button
        self.load_file_button = QPushButton("Load File"); self.load_file_button.setMinimumHeight(BUTTON_MIN_HEIGHT); self.load_file_button.setMinimumWidth(LOAD_FILE_BUTTON_MIN_WIDTH); self.load_file_button.clicked.connect(self.load_file); button_layout.addWidget(self.load_file_button)
        # Process Button
        self.process_button = QPushButton("Process with AI & Generate PDF"); self.process_button.setMinimumHeight(BUTTON_MIN_HEIGHT); self.process_button.setMinimumWidth(PROCESS_BUTTON_MIN_WIDTH); self.process_button.clicked.connect(self.start_ai_processing); button_layout.addWidget(self.process_button)
        page_layout.addWidget(button_container, alignment=Qt.AlignmentFlag.AlignBottom) # Align whole group bottom
        # --- End Bottom Buttons ---
        self._load_selected_prompt(self.prompt_combo.currentIndex()) 
        self._request_token_update() 
        return page
    def _create_status_page(self):
        page = QWidget(); page_layout = QVBoxLayout(page); page_layout.setContentsMargins(0, 0, 0, 0); page_layout.setSpacing(LAYOUT_SPACING); page_layout.addWidget(self.status_label); page_layout.addWidget(self.status_display); page_layout.addWidget(self.log_label); page_layout.addWidget(self.log_display); page_layout.addSpacerItem(QSpacerItem(20, 20, QSizePolicy.Policy.Minimum, QSizePolicy.Policy.Expanding)); back_button_container = QWidget(); back_button_layout = QHBoxLayout(back_button_container); back_button_layout.setContentsMargins(0,0,0,0); back_button_layout.addStretch(1); self.back_button = QPushButton("Back to Input"); self.back_button.setMinimumHeight(BUTTON_MIN_HEIGHT); self.back_button.setMinimumWidth(LOAD_FILE_BUTTON_MIN_WIDTH); self.back_button.clicked.connect(self._go_to_input_page); back_button_layout.addWidget(self.back_button); page_layout.addWidget(back_button_container, alignment=Qt.AlignmentFlag.AlignRight); return page
    def _load_selected_prompt(self, index):
        if hasattr(self, 'prompt_input') and self.prompt_input and self.prompt_input.isEnabled(): selected_prompt_name = self.prompt_combo.itemText(index); prompt_text = PREDEFINED_PROMPTS.get(selected_prompt_name, ""); self.prompt_input.setText(prompt_text); 
        if hasattr(self, 'log_message'): self.log_message(f"Loaded prompt: '{selected_prompt_name}'.")
    def _get_tokenizer(self):
        if self._tokenizer_instance is None:
            try:
                self.log_message("Loading tokenizer...", color=COLOR_STATUS_DEFAULT); QApplication.processEvents() 
                import logging; logging.getLogger("transformers").setLevel(logging.ERROR)
                self._tokenizer_instance = AutoTokenizer.from_pretrained("gpt2")
                self.log_message("Tokenizer loaded.", color=COLOR_STATUS_DEFAULT)
            except Exception as e:
                self._tokenizer_instance = None 
                self.log_message(f"Error loading tokenizer: {e}", color=COLOR_ERROR_RED)
                QMessageBox.warning(self, "Tokenizer Error", f"Could not load tokenizer: {e}\nToken counting may be inaccurate.")
        return self._tokenizer_instance
    def _get_token_count(self, text):
        tokenizer = self._get_tokenizer()
        if tokenizer:
            try: token_ids = tokenizer.encode(text, add_special_tokens=False); return len(token_ids)
            except Exception as e:
                if not hasattr(self, '_tokenizer_error_logged'): self.log_message(f"Error during token counting: {e}", color=COLOR_ERROR_RED); self._tokenizer_error_logged = True
                return 0 
        else: return len(text.split()) if text else 0
    def _request_token_update(self):
        self.token_update_timer.start() 
    def _perform_token_update(self):
        if not hasattr(self, 'original_text_input') or not hasattr(self, 'token_count_label'): return 
        input_text = self.original_text_input.toPlainText(); token_count = self._get_token_count(input_text)
        self.token_count_label.setText(f"Tokens: {token_count} / {LLM_CONTEXT_WINDOW}")
        base_style = "font-family: 'Consolas', 'Monaco', 'Courier New', monospace; font-size: 10px;"
        if token_count > LLM_CONTEXT_WINDOW: self.token_count_label.setStyleSheet(f"color: {COLOR_TOKEN_EXCEEDED}; {base_style}")
        elif token_count > LLM_CONTEXT_WINDOW * (TOKEN_WARNING_THRESHOLD_PERCENT / 100.0): self.token_count_label.setStyleSheet(f"color: {COLOR_TOKEN_WARNING}; {base_style}")
        else: self.token_count_label.setStyleSheet(f"color: {COLOR_TOKEN_NORMAL}; {base_style}")
    def _get_stylesheet(self): 
         return f""" QMainWindow {{ color: {COLOR_TEXT_NEON_GREEN}; font-family: 'Consolas', 'Monaco', 'Courier New', monospace; background-color: {COLOR_BACKGROUND_DARK}; }} QWidget#centralWidget {{ background-color: transparent; }} QWidget {{ background-color: transparent; font-family: 'Consolas', 'Monaco', 'Courier New', monospace; }} QFrame {{ border: 1px solid {COLOR_HIGHLIGHT_CYAN}; border-radius: 5px; background-color: #222222; }} QFrame[frameShape="4"] {{ border: none; background-color: transparent; min-height: {SEPARATOR_HEIGHT_PX}px; max-height: {SEPARATOR_HEIGHT_PX}px; margin-top: {SEPARATOR_MARGIN_V}px; margin-bottom: {SEPARATOR_MARGIN_V}px; background: qlineargradient(x1:0, y1:0, x2:1, y2:0, stop:0 transparent, stop:0.2 {SEPARATOR_COLOR}, stop:0.8 {SEPARATOR_COLOR}, stop:1 transparent); }} QLabel {{ color: {COLOR_HIGHLIGHT_CYAN}; font-size: 14px; font-weight: bold; margin-bottom: 3px; background-color: transparent; }} QLabel:disabled {{ color: {COLOR_DISABLED_TEXT}; }} QTextEdit, QLineEdit, QSpinBox, QComboBox {{ background-color: {COLOR_INPUT_BACKGROUND}; color: {COLOR_TEXT_NEON_GREEN}; border: 1px solid {COLOR_HIGHLIGHT_CYAN}; padding: 8px; selection-background-color: #004444; font-size: 13px; border-radius: 4px; }} QTextEdit:disabled, QLineEdit:disabled, QSpinBox:disabled, QComboBox:disabled {{ background-color: {COLOR_DISABLED_BG}; color: {COLOR_DISABLED_TEXT}; border: 1px solid {COLOR_DISABLED_BORDER}; }} QComboBox::drop-down {{ border-left: 1px solid {COLOR_HIGHLIGHT_CYAN}; }} QComboBox::drop-down:disabled {{ border-left: 1px solid {COLOR_DISABLED_BORDER}; }} QComboBox QAbstractItemView {{ background-color: {COLOR_INPUT_BACKGROUND}; color: {COLOR_TEXT_NEON_GREEN}; selection-background-color: #005555; border: 1px solid {COLOR_HIGHLIGHT_CYAN}; }} QLineEdit {{ padding: 5px 8px; }} QSpinBox {{ padding: 5px 8px; }} QPushButton {{ background-color: {COLOR_BUTTON_GREEN_DARK}; color: {COLOR_TEXT_NEON_GREEN}; border: 2px solid {COLOR_TEXT_NEON_GREEN}; padding: 10px 15px; font-size: 14px; font-weight: bold; border-radius: 7px; }} QPushButton:hover {{ background-color: {COLOR_BUTTON_GREEN_HOVER}; border-color: #33ff33; }} QPushButton:pressed {{ background-color: {COLOR_BUTTON_GREEN_PRESSED}; }} QPushButton:disabled {{ background-color: {COLOR_DISABLED_BG}; color: {COLOR_DISABLED_TEXT}; border-color: {COLOR_DISABLED_BORDER}; }} QMessageBox {{ background-color: {COLOR_BACKGROUND_DARK}; }} QMessageBox QLabel {{ color: {COLOR_TEXT_NEON_GREEN}; font-size: 13px; }} QFileDialog {{ background-color: {COLOR_BACKGROUND_DARK}; }} QProgressDialog {{ background-color: {COLOR_BACKGROUND_DARK}; color: {COLOR_TEXT_NEON_GREEN}; border: 2px solid {COLOR_HIGHLIGHT_CYAN}; border-radius: 5px; }} QProgressDialog QLabel {{ color: {COLOR_HIGHLIGHT_CYAN}; font-size: 13px; font-weight: normal; }} QProgressDialog QProgressBar {{ border: 1px solid {COLOR_TEXT_NEON_GREEN}; border-radius: 3px; background-color: #333333; text-align: center; color: {COLOR_TEXT_NEON_GREEN}; }} QProgressDialog QProgressBar::chunk {{ background-color: {COLOR_TEXT_NEON_GREEN}; }} QProgressDialog QPushButton {{ padding: 5px 10px; font-size: 12px; }} """
    def _apply_background_image(self):
        central_widget = self.centralWidget();
        if not central_widget or not os.path.exists(BACKGROUND_IMAGE_PATH):
            if not os.path.exists(BACKGROUND_IMAGE_PATH) and hasattr(self, 'log_message'): self.log_message(f"Warning: BG image missing: {BACKGROUND_IMAGE_PATH}", COLOR_WARNING_YELLOW)
            return
        try:
            background_pixmap = QPixmap(BACKGROUND_IMAGE_PATH)
            if not background_pixmap.isNull():
                palette = central_widget.palette(); scaled_pixmap = background_pixmap.scaled(central_widget.size(), Qt.AspectRatioMode.IgnoreAspectRatio, Qt.TransformationMode.SmoothTransformation)
                brush = QBrush(scaled_pixmap); palette.setBrush(QPalette.ColorRole.Window, brush); central_widget.setPalette(palette); central_widget.setAutoFillBackground(True)
            elif hasattr(self, 'log_message'): self.log_message(f"Warning: Could not load pixmap from {BACKGROUND_IMAGE_PATH}.", COLOR_WARNING_YELLOW)
        except Exception as e:
            if hasattr(self, 'log_message'): self.log_message(f"Error applying background: {e}", color=COLOR_ERROR_RED)
    def resizeEvent(self, event):
        super().resizeEvent(event); self._apply_background_image()

    def _extract_text_from_docx(self, file_path):
        if docx is None: self.log_message("python-docx not available.", COLOR_ERROR_RED); return None, "python-docx library not installed."
        try:
            doc_obj = docx.Document(file_path); output_lines = [] 
            for para_idx, para in enumerate(doc_obj.paragraphs):
                paragraph_text_parts = []
                for run in para.runs:
                    run_text = run.text
                    if run_text: 
                        run_text = run_text.replace('&', '&').replace('<', '<').replace('>', '>')
                        prefix = ""; suffix = ""
                        if run.bold: prefix += "<b>"; suffix = "</b>" + suffix
                        if run.italic: prefix += "<i>"; suffix = "</i>" + suffix
                        if WD_UNDERLINE and run.underline and run.underline != WD_UNDERLINE.NONE and not run.strike: prefix += "<u>"; suffix = "</u>" + suffix
                        elif not WD_UNDERLINE and run.underline and not run.strike: prefix += "<u>"; suffix = "</u>" + suffix
                        paragraph_text_parts.append(prefix + run_text + suffix)
                full_paragraph_text = "".join(paragraph_text_parts).strip()
                style_name = para.style.name if para.style and para.style.name else ""; style_name_lower = style_name.lower()
                is_list_style_by_name = 'list paragraph' in style_name_lower or 'list bullet' in style_name_lower or 'list number' in style_name_lower
                is_list_by_text_pattern = bool(re.match(r"^\s*(\*|-|•|▪|o)\s+", full_paragraph_text)) or bool(re.match(r"^\s*(\d+\.|[a-zA-Z][\.\)])\s+", full_paragraph_text))
                line_added = False
                if style_name.startswith('Heading 1'):
                    if para_idx > 0 and output_lines and output_lines[-1] != "": output_lines.append("") 
                    output_lines.append(f"# {full_paragraph_text}"); line_added = True
                elif style_name.startswith('Heading 2'):
                    if para_idx > 0 and output_lines and output_lines[-1] != "": output_lines.append("")
                    output_lines.append(f"## {full_paragraph_text}"); line_added = True
                elif style_name.startswith('Heading 3'):
                    if para_idx > 0 and output_lines and output_lines[-1] != "": output_lines.append("")
                    output_lines.append(f"### {full_paragraph_text}"); line_added = True
                elif not line_added and (is_list_style_by_name or is_list_by_text_pattern):
                    if full_paragraph_text: 
                        text_without_prefix = re.sub(r"^\s*(\*|-|•|▪|o|\d+\.|[a-zA-Z][\.\)])\s+", "", full_paragraph_text).strip()
                        output_lines.append(f"* {text_without_prefix}"); line_added = True
                if not line_added:
                    if full_paragraph_text: output_lines.append(full_paragraph_text)
                    elif para_idx > 0 and output_lines and output_lines[-1] != "": output_lines.append("")
            final_text = "\n".join(output_lines)
            final_text = re.sub(r'\n(\s*\n)+', '\n\n', final_text).strip() 
            return final_text, None
        except Exception as e:
            self.log_message(f"Error reading .docx file {os.path.basename(file_path)}: {e}", COLOR_ERROR_RED)
            return None, f"Error reading .docx file: {e}"

    # +++ ADDED Method to Extract Text from PPTX +++
    def _extract_text_from_pptx(self, file_path):
        """Extracts text from a PPTX file slide by slide."""
        if Presentation is None:
            self.log_message("python-pptx library is not available. Cannot process PPTX files.", COLOR_ERROR_RED)
            return None, "python-pptx library not installed."
        
        full_text = []
        try:
            prs = Presentation(file_path)
            for i, slide in enumerate(prs.slides):
                slide_texts = []
                
                # Add notes first
                if slide.has_notes_slide:
                    notes_frame = slide.notes_slide.notes_text_frame
                    if notes_frame and notes_frame.text and notes_frame.text.strip():
                        notes_text = notes_frame.text.strip()
                        if slide_texts or full_text: slide_texts.append("") 
                        slide_texts.append(f"### Speaker Notes (Slide {i+1}):")
                        for note_para in notes_text.split('\n'):
                            if note_para.strip():
                                slide_texts.append(html.escape(note_para.strip())) # Escape notes text
                        slide_texts.append("") 
                
                # Add slide title 
                title_text = ""
                if slide.shapes.title and slide.shapes.title.has_text_frame and slide.shapes.title.text.strip():
                     title_text = html.escape(slide.shapes.title.text.strip())
                     if slide_texts and slide_texts[-1] != "": slide_texts.append("") 
                     slide_texts.append(f"## Slide {i+1}: {title_text}") 

                # Extract text from other shapes
                shape_texts = []
                for shape in slide.shapes:
                    if not shape.has_text_frame: continue
                    if shape == slide.shapes.title: continue 
                    if slide.has_notes_slide and shape.is_placeholder and shape.name.startswith("Notes Placeholder"): continue

                    shape_content = []
                    for paragraph in shape.text_frame.paragraphs:
                        para_text_parts = []
                        for run in paragraph.runs:
                            run_text = run.text 
                            if run_text:
                                escaped_run_text = html.escape(run_text) # Escape text content
                                prefix = ""; suffix = ""
                                if run.font.bold: prefix += "<b>"; suffix = "</b>" + suffix
                                if run.font.italic: prefix += "<i>"; suffix = "</i>" + suffix
                                # Simple underline check (needs WD_UNDERLINE imported)
                                if WD_UNDERLINE and run.font.underline and run.font.underline != WD_UNDERLINE.NONE: prefix += "<u>"; suffix = "</u>" + suffix
                                elif not WD_UNDERLINE and run.font.underline: prefix += "<u>"; suffix = "</u>" + suffix
                                para_text_parts.append(prefix + escaped_run_text + suffix)
                        
                        para_full_text = "".join(para_text_parts).strip()
                        if para_full_text:
                            indent_level = paragraph.level
                            # Prefix list items based on indentation level
                            if indent_level > 0:
                                shape_content.append(f"* {para_full_text}") 
                            else:
                                shape_content.append(para_full_text)
                    
                    if shape_content:
                        # Join paragraphs within a shape with single newlines
                        shape_texts.append("\n".join(shape_content)) 

                if shape_texts:
                     # Add space before shape text if title or notes exist
                     if slide_texts and slide_texts[-1] != "": slide_texts.append("") 
                     # Join different shapes with double newlines
                     slide_texts.append("\n\n".join(shape_texts)) 

                if slide_texts:
                    full_text.append("\n".join(slide_texts)) 

            final_text = "\n\n".join(full_text)
            final_text = re.sub(r'\n(\s*\n)+', '\n\n', final_text).strip()
            return final_text, None

        except Exception as e:
            error_msg = f"Error reading PPTX file {os.path.basename(file_path)}: {e}"
            self.log_message(error_msg, color=COLOR_ERROR_RED)
            # import traceback # Uncomment for detailed debug
            # self.log_message(traceback.format_exc(), color=COLOR_ERROR_RED) # Uncomment for detailed debug
            return None, f"Error reading PPTX file: {e}"


    def load_file(self):
        # --- MODIFIED to include PPTX ---
        file_types = "Supported Files (*.txt *.docx *.pptx);;Text files (*.txt);;Word documents (*.docx);;PowerPoint files (*.pptx);;All files (*)" # Added PPTX filter
        file_path, selected_filter = QFileDialog.getOpenFileName(
            self, "Select a File to Load", "", file_types)

        if file_path:
            file_content = None; error_message = None
            file_extension = os.path.splitext(file_path)[1].lower()

            try:
                if file_extension == '.txt':
                    with open(file_path, 'r', encoding='utf-8', errors='ignore') as f: file_content = f.read()
                elif file_extension == '.docx':
                    if docx is None: error_message = "python-docx library not installed."; QMessageBox.warning(self, "Library Missing", error_message)
                    else: file_content, error_message = self._extract_text_from_docx(file_path)
                # +++ Add PPTX handling +++
                elif file_extension == '.pptx':
                    if Presentation is None: # Check if Presentation class is available
                        error_message = "python-pptx library not installed. Cannot load .pptx files."
                        QMessageBox.warning(self, "Library Missing", error_message)
                    else:
                        file_content, error_message = self._extract_text_from_pptx(file_path) # Call the extractor
                # +++ End PPTX handling +++
                else: 
                    error_message = f"Unsupported file type: {file_extension}. Please select .txt, .docx, or .pptx." # Updated message
                    QMessageBox.warning(self, "Unsupported File", error_message)
                
                # Process result
                if file_content is not None:
                    self.original_text_input.setText(file_content); self.update_status(f"Loaded file: {os.path.basename(file_path)}", COLOR_TEXT_NEON_GREEN)
                    self.log_message(f"Loaded file: {os.path.basename(file_path)}"); self._request_token_update() 
                elif error_message: 
                    self.update_status(f"Error loading {os.path.basename(file_path)}.", COLOR_ERROR_RED); self.log_message(error_message, color=COLOR_ERROR_RED) 
                    if "Unsupported file type" not in error_message and "library not installed" not in error_message: QMessageBox.critical(self, "Error", error_message)
                
            except Exception as e: 
                error_msg = f"Error processing file {os.path.basename(file_path)}: {e}"; self.update_status("Error processing file.", COLOR_ERROR_RED)
                self.log_message(error_msg, color=COLOR_ERROR_RED); QMessageBox.critical(self, "Error", error_msg)


    # --- Subsequent methods (start_ai_processing, cancel_ai_processing, etc.) are correct ---
    def start_ai_processing(self): 
        original_text = self.original_text_input.toPlainText().strip(); prompt_instruction = self.prompt_input.toPlainText().strip() 
        if not original_text: QMessageBox.warning(self, "Input Required", "Please enter or load text."); self.update_status("Please enter or load text.", COLOR_WARNING_YELLOW); self.log_message("Processing cancelled: No original text.", color=COLOR_WARNING_YELLOW); return
        if not prompt_instruction: QMessageBox.warning(self, "Input Required", "Please provide AI instructions."); self.update_status("Please provide AI instructions.", COLOR_WARNING_YELLOW); self.log_message("Processing cancelled: No AI instructions.", color=COLOR_WARNING_YELLOW); return
        self.stacked_widget.setCurrentIndex(1); self.log_display.clear(); self.update_status("Starting AI processing...", COLOR_WARNING_YELLOW); self.log_message("Starting AI processing...")
        self.progress_dialog = QProgressDialog("AI Processing...", "Cancel", 0, 0, self); self.progress_dialog.setWindowTitle("AI at Work"); self.progress_dialog.setWindowModality(Qt.WindowModality.WindowModal); self.progress_dialog.canceled.connect(self.cancel_ai_processing); self.progress_dialog.show()
        self.ai_worker = AIWorker(original_text, prompt_instruction); self.ai_worker.finished.connect(self.handle_ai_response); self.ai_worker.progress.connect(self.update_progress_dialog); self.ai_worker.start()
    def cancel_ai_processing(self):
        self.log_message("AI processing cancellation requested.", color=COLOR_WARNING_YELLOW) 
        if self.ai_worker and self.ai_worker.isRunning(): self.ai_worker.stop()
        if self.progress_dialog: self.progress_dialog.close(); self.progress_dialog = None 
        self.update_status("AI processing cancellation requested.", COLOR_WARNING_YELLOW)
        if hasattr(self, 'back_button') and self.back_button: self.back_button.setEnabled(True)
    def update_progress_dialog(self, message):
         if self.progress_dialog and self.progress_dialog.isVisible(): self.progress_dialog.setLabelText(message)
         self.log_message(f"AI Progress: {message}", color=COLOR_STATUS_DEFAULT)
    def handle_ai_response(self, success, ai_output_or_error, was_cancelled): 
        if self.progress_dialog: self.progress_dialog.close(); self.progress_dialog = None
        if was_cancelled:
            self.log_message("AI processing was confirmed cancelled. Skipping PDF generation.", COLOR_WARNING_YELLOW); self.update_status("AI Processing Cancelled.", COLOR_WARNING_YELLOW)
            if hasattr(self, 'back_button') and self.back_button: self.back_button.setEnabled(True)
            return 
        if not success: 
            self.update_status(f"AI Processing Failed: {ai_output_or_error}", COLOR_ERROR_RED); self.log_message(f"AI Error: {ai_output_or_error}", color=COLOR_ERROR_RED)
            QMessageBox.critical(self, "AI Error", ai_output_or_error)
            if hasattr(self, 'back_button') and self.back_button: self.back_button.setEnabled(True)
            return 
        self.update_status("AI processing complete. Preparing PDF...", COLOR_TEXT_NEON_GREEN); self.log_message("AI processing complete. Preparing PDF generation...")
        self.save_pdf_from_ai_output(ai_output_or_error) 
    def save_pdf_from_ai_output(self, processed_text): 
         selected_page_size = self.page_size_combo.currentText(); selected_font_size = self.font_size_spinbox.value(); default_filename = "ai_formatted_document.pdf" 
         output_filename, _ = QFileDialog.getSaveFileName(self, "Save AI Formatted PDF", default_filename, "PDF files (*.pdf);;All files (*)") 
         if not output_filename: 
            self.update_status("PDF save cancelled.", COLOR_WARNING_YELLOW); self.log_message("PDF save cancelled.", color=COLOR_WARNING_YELLOW); 
            if hasattr(self, 'back_button'): self.back_button.setEnabled(True); 
            print("DEBUG UI: PDF save dialog cancelled by user.") 
            return
         print(f"DEBUG UI: PDF Filename selected: {output_filename}") 
         if not output_filename.lower().endswith('.pdf'): output_filename += '.pdf'
         self.update_status("Generating PDF in background...", COLOR_WARNING_YELLOW); 
         self.log_message(f"Starting background PDF generation '{os.path.basename(output_filename)}'...")
         if hasattr(self, 'back_button'): self.back_button.setEnabled(False)
         print(f"DEBUG UI: Creating PDFWorker (Font: {selected_font_size}). Text length: {len(processed_text)}") 
         self.pdf_worker = PDFWorker(processed_text, output_filename, selected_page_size, selected_font_size)
         print(f"DEBUG UI: Connecting PDFWorker finished signal...") 
         self.pdf_worker.finished.connect(self.handle_pdf_result) 
         print(f"DEBUG UI: Starting PDFWorker...") 
         self.pdf_worker.start()
    def handle_pdf_result(self, success, message):
        print(f"DEBUG UI: handle_pdf_result received: success={success}, message='{message[:100]}...'") 
        if hasattr(self, 'back_button'): self.back_button.setEnabled(True)
        if success: 
            self.update_status(message, COLOR_TEXT_NEON_GREEN); self.log_message(message, color=COLOR_TEXT_NEON_GREEN)
        else: 
            detailed_error_message = f"PDF Generation Error: {message}"
            self.update_status("PDF Generation Failed.", COLOR_ERROR_RED) 
            self.log_message(detailed_error_message, color=COLOR_ERROR_RED) 
            QMessageBox.critical(self, "PDF Generation Error", f"Failed to generate PDF.\nDetails logged.\nError: {message.splitlines()[0]}") 
        self.pdf_worker = None 
    def update_status(self, message, color=COLOR_STATUS_DEFAULT):
        if hasattr(self, 'status_display') and self.status_display: self.status_display.setText(message); self.status_display.setStyleSheet(f"QLineEdit {{ color: {color}; background-color: {COLOR_INPUT_BACKGROUND}; border: 1px solid {COLOR_HIGHLIGHT_CYAN}; padding: 5px 8px; }}")
    def log_message(self, message, color=None): 
        if hasattr(self, 'log_display') and self.log_display:
             timestamp = time.strftime("%H:%M:%S"); log_color = color if color else COLOR_HIGHLIGHT_CYAN 
             current_html = self.log_display.toHtml()
             if current_html.endswith("</body></html>"): current_html = current_html[:-len("</body></html>")]
             escaped_message = html.escape(str(message)) 
             new_log_entry = f"<br><span style='color: {log_color}; font-family: \"Consolas\", \"Monaco\", \"Courier New\", monospace;'>[{timestamp}] {escaped_message}</span>"
             self.log_display.setHtml(current_html + new_log_entry + "</body></html>"); self.log_display.verticalScrollBar().setValue(self.log_display.verticalScrollBar().maximum())
    def _go_to_input_page(self):
        self.stacked_widget.setCurrentIndex(0)
        if hasattr(self, 'process_button') and self.process_button: self.process_button.setEnabled(True) 
        self.update_status("Ready for input.", COLOR_STATUS_DEFAULT)

# --- Standalone Execution ---
if __name__ == '__main__':
    print("Running UI standalone test...")
    app = QApplication(sys.argv)
    main_window = ModernHackerPDFConverterWindow()
    main_window.show()
    sys.exit(app.exec())